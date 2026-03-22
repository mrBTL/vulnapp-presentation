"""
Jules Bot — Telegram Automation
Wizualna prezentacja architektury i przepływów.
Google Slides compatible (no add_connector).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Paleta ──────────────────────────────────────────────────────────
BG          = RGBColor(0x0D, 0x17, 0x26)
CARD        = RGBColor(0x0F, 0x28, 0x44)
C_TELEGRAM  = RGBColor(0x26, 0x9F, 0xDA)   # Telegram blue
C_CLAUDE    = RGBColor(0xCC, 0x77, 0xFF)   # Claude purple
C_GITHUB    = RGBColor(0xFF, 0xFF, 0xFF)   # GitHub white
C_SYSTEM    = RGBColor(0x00, 0xE6, 0x76)   # systemd / server green
C_FILE      = RGBColor(0xFF, 0xBF, 0x00)   # state files amber
C_DECISION  = RGBColor(0xFF, 0x60, 0x60)   # decision red-orange
C_EXEC      = RGBColor(0xFF, 0x45, 0x45)   # execution red
C_ACCENT    = RGBColor(0x00, 0xE5, 0xFF)   # general cyan
TEXT_W      = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DIM    = RGBColor(0x7A, 0x9B, 0xBF)
TEXT_DARK   = RGBColor(0x08, 0x12, 0x1E)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ════════════════════════════════════════════════════════════════════
# PRIMITIVE HELPERS
# ════════════════════════════════════════════════════════════════════

def bg(slide, color=BG):
    r = slide.shapes.add_shape(1, 0, 0, W, H)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background()

def rect(slide, x, y, w, h, fill, line=None, lw=Pt(1.5)):
    r = slide.shapes.add_shape(1, x, y, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = fill
    if line: r.line.color.rgb = line; r.line.width = lw
    else: r.line.fill.background()
    return r

def rrect(slide, x, y, w, h, fill, text="", fs=Pt(12), fc=TEXT_W,
          bold=False, line=None, lw=Pt(1.5), align=PP_ALIGN.CENTER):
    r = slide.shapes.add_shape(5, x, y, w, h)
    r.adjustments[0] = 0.06
    r.fill.solid(); r.fill.fore_color.rgb = fill
    if line: r.line.color.rgb = line; r.line.width = lw
    else: r.line.fill.background()
    if text:
        tf = r.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = align
        ru = p.add_run(); ru.text = text
        ru.font.size = fs; ru.font.bold = bold
        ru.font.color.rgb = fc; ru.font.name = "Calibri"
    return r

def diamond(slide, x, y, w, h, fill, text="", fs=Pt(11), fc=TEXT_W, bold=False, line=None):
    r = slide.shapes.add_shape(4, x, y, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = fill
    if line: r.line.color.rgb = line; r.line.width = Pt(1.5)
    else: r.line.fill.background()
    if text:
        tf = r.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        ru = p.add_run(); ru.text = text
        ru.font.size = fs; ru.font.bold = bold
        ru.font.color.rgb = fc; ru.font.name = "Calibri"
    return r

def oval(slide, x, y, w, h, fill, text="", fs=Pt(12), fc=TEXT_W, bold=False):
    r = slide.shapes.add_shape(9, x, y, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = fill
    r.line.fill.background()
    if text:
        tf = r.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        ru = p.add_run(); ru.text = text
        ru.font.size = fs; ru.font.bold = bold
        ru.font.color.rgb = fc; ru.font.name = "Calibri"
    return r

def txt(slide, text, x, y, w, h, fs=Pt(13), bold=False, color=TEXT_W,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    ru = p.add_run(); ru.text = text
    ru.font.size = fs; ru.font.bold = bold; ru.font.italic = italic
    ru.font.color.rgb = color; ru.font.name = "Calibri"
    return tb

def hline(slide, x1, x2, cy, color, thick=0.035):
    rect(slide, x1, cy - Inches(thick/2), x2 - x1, Inches(thick), color)

def vline(slide, cx, y1, y2, color, thick=0.035):
    rect(slide, cx - Inches(thick/2), y1, Inches(thick), y2 - y1, color)

def arrow_d(slide, cx, y1, y2, color):
    """Arrow pointing down."""
    vline(slide, cx, y1, y2 - Inches(0.10), color)
    tw = Inches(0.16)
    rect(slide, cx - tw/2,           y2 - Inches(0.10), tw,          Inches(0.05), color)
    rect(slide, cx - tw*0.6/2,       y2 - Inches(0.05), tw * 0.6,   Inches(0.04), color)
    rect(slide, cx - tw*0.2/2,       y2 - Inches(0.01), tw * 0.2,   Inches(0.03), color)

def arrow_r(slide, x1, x2, cy, color):
    """Arrow pointing right."""
    hline(slide, x1, x2 - Inches(0.10), cy, color)
    th = Inches(0.16)
    rect(slide, x2 - Inches(0.10), cy - th/2,           Inches(0.05), th,         color)
    rect(slide, x2 - Inches(0.05), cy - th*0.6/2,       Inches(0.04), th * 0.6,   color)
    rect(slide, x2 - Inches(0.01), cy - th*0.2/2,       Inches(0.03), th * 0.2,   color)

def arrow_l(slide, x1, x2, cy, color):
    """Arrow pointing left (x1 > x2)."""
    hline(slide, x2 + Inches(0.10), x1, cy, color)
    th = Inches(0.16)
    rect(slide, x2 + Inches(0.05), cy - th/2,         Inches(0.05), th,         color)
    rect(slide, x2 + Inches(0.01), cy - th*0.6/2,     Inches(0.04), th * 0.6,   color)
    rect(slide, x2,                cy - th*0.2/2,     Inches(0.03), th * 0.2,   color)

def arrow_u(slide, cx, y1, y2, color):
    """Arrow pointing up (y1 > y2)."""
    vline(slide, cx, y2 + Inches(0.10), y1, color)
    tw = Inches(0.16)
    rect(slide, cx - tw/2,       y2 + Inches(0.05), tw,        Inches(0.05), color)
    rect(slide, cx - tw*0.6/2,   y2 + Inches(0.01), tw*0.6,   Inches(0.04), color)
    rect(slide, cx - tw*0.2/2,   y2,                tw*0.2,   Inches(0.03), color)

def label_on_line(slide, text, x, y, color=TEXT_DIM):
    txt(slide, text, x, y, Inches(0.8), Inches(0.28),
        fs=Pt(9), color=color, bold=True, align=PP_ALIGN.CENTER)

def slide_num(slide, n, total=5):
    txt(slide, f"0{n} / 0{total}", Inches(11.9), Inches(7.0), Inches(1.3), Inches(0.35),
        fs=Pt(10), color=C_ACCENT, align=PP_ALIGN.RIGHT)

def header(slide, title, subtitle, accent_color=C_ACCENT):
    rect(slide, 0, 0, W, Inches(1.3), RGBColor(0x07, 0x14, 0x22))
    txt(slide, title, Inches(0.5), Inches(0.12), Inches(11), Inches(0.65),
        fs=Pt(30), bold=True, color=TEXT_W)
    txt(slide, subtitle, Inches(0.5), Inches(0.75), Inches(11), Inches(0.4),
        fs=Pt(13), color=TEXT_DIM, italic=True)
    rect(slide, Inches(0.5), Inches(1.25), Inches(12.33), Inches(0.04), accent_color)


# ════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════

s1 = prs.slides.add_slide(BLANK)
bg(s1)

# Left dark panel
rect(s1, 0, 0, Inches(6.0), H, RGBColor(0x06, 0x12, 0x22))

# Title text
txt(s1, "Jules Bot", Inches(0.6), Inches(1.6), Inches(5.2), Inches(1.4),
    fs=Pt(56), bold=True, color=C_TELEGRAM)
txt(s1, "Telegram Automation Layer", Inches(0.6), Inches(3.1), Inches(5.2), Inches(0.6),
    fs=Pt(20), color=TEXT_W)
rect(s1, Inches(0.6), Inches(2.98), Inches(1.4), Inches(0.05), C_TELEGRAM)
txt(s1, "dla VulnApp-Light · rafserver",
    Inches(0.6), Inches(3.75), Inches(5.2), Inches(0.45),
    fs=Pt(14), color=TEXT_DIM, italic=True)

# Right — component circles
components = [
    (Inches(7.2),  Inches(1.5), "jules_listener.sh",  "systemd\nlong-poll",      C_SYSTEM),
    (Inches(10.2), Inches(1.5), "jules_actions.sh",   "brain\n2-phase",           C_TELEGRAM),
    (Inches(7.2),  Inches(4.0), "jules_review.sh",    "cron\n5 min",              C_CLAUDE),
    (Inches(10.2), Inches(4.0), "debate.sh",          "Claude ↔\nGemini",         C_FILE),
]
for cx, cy, title, sub, col in components:
    oval(s1, cx - Inches(0.9), cy - Inches(0.75), Inches(1.8), Inches(1.5),
         RGBColor(int(col[0]*0.3), int(col[1]*0.3), int(col[2]*0.3)),
         bold=False)
    rect(s1, cx - Inches(0.9), cy - Inches(0.75), Inches(1.8), Inches(1.5),
         RGBColor(int(col[0]*0.3), int(col[1]*0.3), int(col[2]*0.3)),
         line=col, lw=Pt(2.0))
    txt(s1, title, cx - Inches(0.88), cy - Inches(0.62), Inches(1.76), Inches(0.4),
        fs=Pt(10), bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s1, sub,   cx - Inches(0.88), cy - Inches(0.18), Inches(1.76), Inches(0.5),
        fs=Pt(10), color=TEXT_DIM, align=PP_ALIGN.CENTER)

# Connection arrows between components
arrow_r(s1, Inches(8.1), Inches(9.3), Inches(1.5 - 0.1), C_SYSTEM)
arrow_d(s1, Inches(7.2), Inches(2.25), Inches(3.25), C_CLAUDE)
arrow_d(s1, Inches(10.2), Inches(2.25), Inches(3.25), C_TELEGRAM)
arrow_r(s1, Inches(8.1), Inches(9.3), Inches(4.0 - 0.1), C_FILE)

# Legend
items = [
    (C_TELEGRAM,  "Telegram API"),
    (C_CLAUDE,    "Claude / Gemini AI"),
    (C_SYSTEM,    "Server / systemd"),
    (C_FILE,      "State files / cron"),
    (C_GITHUB,    "GitHub CLI"),
]
for i, (col, lbl) in enumerate(items):
    bx = Inches(6.4) + i * Inches(1.38)
    rect(s1, bx, Inches(6.85), Inches(0.22), Inches(0.22), col)
    txt(s1, lbl, bx + Inches(0.28), Inches(6.82), Inches(1.1), Inches(0.28),
        fs=Pt(9.5), color=TEXT_DIM)

slide_num(s1, 1)


# ════════════════════════════════════════════════════════════════════
# SLIDE 2 — TOPOLOGY MAP
# ════════════════════════════════════════════════════════════════════

s2 = prs.slides.add_slide(BLANK)
bg(s2)
header(s2, "Mapa Zależności — Topology", "Wszystkie komponenty i połączenia", C_TELEGRAM)

# ── ROW 1: External world ────────────────────────────────────────
# User (iPhone)
rrect(s2, Inches(0.25), Inches(1.55), Inches(1.8), Inches(0.75),
      RGBColor(0x0F, 0x28, 0x44), "📱 User\niPhone", fs=Pt(11),
      bold=True, fc=TEXT_W, line=C_TELEGRAM)

# Telegram API
rrect(s2, Inches(2.6), Inches(1.55), Inches(2.2), Inches(0.75),
      RGBColor(0x06, 0x38, 0x5C), "☁️  Telegram API\napi.telegram.org",
      fs=Pt(10), bold=True, fc=C_TELEGRAM, line=C_TELEGRAM)

# GitHub
rrect(s2, Inches(5.4), Inches(1.55), Inches(2.2), Inches(0.75),
      RGBColor(0x12, 0x12, 0x12), "🐙  GitHub API\nmrBTL/VulnApp-Light",
      fs=Pt(10), bold=True, fc=C_GITHUB, line=C_GITHUB)

# Claude API
rrect(s2, Inches(8.3), Inches(1.55), Inches(2.3), Inches(0.75),
      RGBColor(0x2A, 0x0A, 0x44), "🤖  Claude API\nsonnet-4-6 / haiku",
      fs=Pt(10), bold=True, fc=C_CLAUDE, line=C_CLAUDE)

# Gemini
rrect(s2, Inches(11.1), Inches(1.55), Inches(1.95), Inches(0.75),
      RGBColor(0x0D, 0x2A, 0x1A), "✨  Gemini\ngemini -p",
      fs=Pt(10), bold=True, fc=C_SYSTEM, line=C_SYSTEM)

# ── ROW 2: Services ──────────────────────────────────────────────
# jules_listener.sh (systemd)
rrect(s2, Inches(0.25), Inches(3.1), Inches(2.5), Inches(1.1),
      RGBColor(0x06, 0x2A, 0x1A),
      "⚙️  jules-listener.service\n(systemd, autostart)\nlong-poll timeout=30s",
      fs=Pt(10), bold=False, fc=C_SYSTEM, line=C_SYSTEM)

# jules_actions.sh — central brain
rrect(s2, Inches(3.2), Inches(2.85), Inches(3.6), Inches(1.5),
      RGBColor(0x06, 0x28, 0x50),
      "🧠  jules_actions.sh\nlock: PID + flock\n2-phase: plan → confirm → exec\nhistory · context · routing",
      fs=Pt(10), bold=False, fc=C_TELEGRAM, line=C_TELEGRAM, lw=Pt(2.5))

# State files cluster
rrect(s2, Inches(7.2), Inches(2.9), Inches(0.95), Inches(0.5),
      RGBColor(0x2A, 0x1A, 0x00), ".tg_offset",    fs=Pt(9), fc=C_FILE, line=C_FILE)
rrect(s2, Inches(7.2), Inches(3.5), Inches(0.95), Inches(0.5),
      RGBColor(0x2A, 0x1A, 0x00), ".tg_history",   fs=Pt(9), fc=C_FILE, line=C_FILE)
rrect(s2, Inches(8.3), Inches(2.9), Inches(1.1), Inches(0.5),
      RGBColor(0x2A, 0x1A, 0x00), ".pending_action",fs=Pt(8.5), fc=C_FILE, line=C_FILE)
rrect(s2, Inches(8.3), Inches(3.5), Inches(1.1), Inches(0.5),
      RGBColor(0x2A, 0x1A, 0x00), ".pending_jules", fs=Pt(8.5), fc=C_FILE, line=C_FILE)
rrect(s2, Inches(9.55), Inches(2.9), Inches(1.1), Inches(0.5),
      RGBColor(0x2A, 0x1A, 0x00), ".debate_topic",  fs=Pt(8.5), fc=C_FILE, line=C_FILE)

txt(s2, "STATE FILES", Inches(7.2), Inches(2.62), Inches(3.5), Inches(0.3),
    fs=Pt(9), bold=True, color=C_FILE)

# jules_review.sh
rrect(s2, Inches(0.25), Inches(4.9), Inches(2.5), Inches(0.9),
      RGBColor(0x1A, 0x06, 0x2A),
      "🔄  jules_review.sh\ncron każde 5 min\n1 PR per run",
      fs=Pt(10), fc=C_CLAUDE, line=C_CLAUDE)

# .reviewed_prs
rrect(s2, Inches(3.2), Inches(5.05), Inches(1.1), Inches(0.5),
      RGBColor(0x2A, 0x1A, 0x00), ".reviewed_prs", fs=Pt(9), fc=C_FILE, line=C_FILE)

# debate.sh
rrect(s2, Inches(4.8), Inches(4.9), Inches(2.2), Inches(0.9),
      RGBColor(0x1A, 0x15, 0x00),
      "🗣️  debate.sh\nClaude ↔ Gemini rounds\nmax 6 · consensus check",
      fs=Pt(10), fc=C_FILE, line=C_FILE)

# Jules AI (remote agent)
rrect(s2, Inches(7.2), Inches(4.9), Inches(2.1), Inches(0.9),
      RGBColor(0x10, 0x22, 0x10),
      "🤖  Jules AI\n(remote agent)\njules new --repo",
      fs=Pt(10), fc=C_SYSTEM, line=C_SYSTEM)

# ── ARROWS ──────────────────────────────────────────────────────
# User ↔ Telegram API (double arrow via labels)
arrow_r(s2, Inches(2.05), Inches(2.60), Inches(1.95), C_TELEGRAM)
arrow_l(s2, Inches(2.05), Inches(2.60), Inches(2.18), C_TELEGRAM)

# Telegram API → listener
arrow_d(s2, Inches(2.7), Inches(2.30), Inches(3.10), C_TELEGRAM)
txt(s2, "getUpdates\npoll", Inches(2.72), Inches(2.55), Inches(1.1), Inches(0.5),
    fs=Pt(8), color=C_TELEGRAM)

# listener → actions
arrow_r(s2, Inches(2.75), Inches(3.20), Inches(3.65), C_SYSTEM)
txt(s2, "trigger", Inches(2.8), Inches(3.45), Inches(0.8), Inches(0.25),
    fs=Pt(8), color=C_SYSTEM)

# actions → Telegram (sendMessage)
arrow_u(s2, Inches(3.7), Inches(2.85), Inches(2.30), C_TELEGRAM)
txt(s2, "sendMessage\nbottons", Inches(2.42), Inches(2.52), Inches(1.0), Inches(0.4),
    fs=Pt(8), color=C_TELEGRAM)

# actions → GitHub
arrow_r(s2, Inches(6.8), Inches(5.40), Inches(2.0), C_GITHUB)
txt(s2, "gh cli", Inches(5.95), Inches(1.78), Inches(0.7), Inches(0.25),
    fs=Pt(8), bold=True, color=C_GITHUB)

# actions → Claude (plan)
arrow_r(s2, Inches(6.8), Inches(8.30), Inches(3.25), C_CLAUDE)
txt(s2, "plan\nphase", Inches(7.15), Inches(3.05), Inches(0.8), Inches(0.4),
    fs=Pt(8), color=C_CLAUDE)

# actions → Claude (exec)
arrow_r(s2, Inches(6.8), Inches(8.30), Inches(3.65), C_EXEC)
txt(s2, "exec\n+Bash(*)", Inches(7.1), Inches(3.62), Inches(0.9), Inches(0.4),
    fs=Pt(8), bold=True, color=C_EXEC)

# actions ↔ state files
hline(s2, Inches(6.8), Inches(7.2), Inches(3.25), C_FILE, thick=0.025)
hline(s2, Inches(6.8), Inches(7.2), Inches(3.65), C_FILE, thick=0.025)

# review → .reviewed_prs
arrow_r(s2, Inches(2.75), Inches(3.2), Inches(5.3), C_FILE)

# review → Claude
arrow_r(s2, Inches(2.75), Inches(8.3), Inches(5.35), C_CLAUDE)
txt(s2, "review", Inches(5.5), Inches(5.12), Inches(0.7), Inches(0.25),
    fs=Pt(8), color=C_CLAUDE)

# actions → debate.sh
arrow_d(s2, Inches(5.0), Inches(4.35), Inches(4.90), C_FILE)
txt(s2, "/debate", Inches(5.05), Inches(4.57), Inches(0.7), Inches(0.25),
    fs=Pt(8), color=C_FILE)

# debate → Claude
arrow_r(s2, Inches(7.0), Inches(8.3), Inches(5.17), C_CLAUDE)
# debate → Gemini
arrow_r(s2, Inches(7.0), Inches(11.1), Inches(5.38), C_SYSTEM)
txt(s2, "rounds", Inches(9.0), Inches(5.18), Inches(0.8), Inches(0.25),
    fs=Pt(8), color=C_CLAUDE)

# actions → Jules AI (delegation)
arrow_r(s2, Inches(6.8), Inches(7.2), Inches(5.3), C_SYSTEM)
txt(s2, "delegate", Inches(6.82), Inches(5.08), Inches(0.8), Inches(0.25),
    fs=Pt(8), color=C_SYSTEM)

# Claude ↔ Gemini (in debate)
arrow_r(s2, Inches(10.6), Inches(11.1), Inches(1.93), C_CLAUDE)
arrow_l(s2, Inches(10.6), Inches(11.1), Inches(2.19), C_SYSTEM)

slide_num(s2, 2)


# ════════════════════════════════════════════════════════════════════
# SLIDE 3 — MAIN MESSAGE FLOW (State Machine)
# ════════════════════════════════════════════════════════════════════

s3 = prs.slides.add_slide(BLANK)
bg(s3)
header(s3, "Przepływ Wiadomości — State Machine", "Od Telegram → Plan → Confirm → Exec", C_TELEGRAM)

# Two columns: left = guards/routing, right = phase detail
# Left column centers
LC = Inches(2.2)   # left column center x
RC = Inches(8.8)   # right column center x
NW = Inches(3.2)   # node width
NH = Inches(0.62)  # node height
DW = Inches(3.0)   # diamond width
DH = Inches(0.75)  # diamond height

# ── LEFT COLUMN: Main flow ──────────────────────────────────────
y = Inches(1.5)

# START
oval(s3, LC - Inches(1.1), y, Inches(2.2), Inches(0.5),
     C_SYSTEM, "▶  Wiadomość z Telegram", fs=Pt(11), bold=True)
y += Inches(0.5)
arrow_d(s3, LC, y, y + Inches(0.3), C_SYSTEM)
y += Inches(0.3)

# Lock check
diamond(s3, LC - DW/2, y, DW, DH, RGBColor(0x2A, 0x10, 0x00),
        "PID istnieje\ni działa?", fs=Pt(11), bold=True,
        fc=C_DECISION, line=C_DECISION)
# YES branch → exit
arrow_r(s3, LC + DW/2, LC + DW/2 + Inches(0.8), y + DH/2, C_DECISION)
rrect(s3, LC + DW/2 + Inches(0.8), y + DH/2 - Inches(0.22),
      Inches(1.3), Inches(0.44),
      RGBColor(0x28, 0x06, 0x06), "⛔  exit(0)\nalready running",
      fs=Pt(9), fc=C_DECISION, line=C_DECISION)
txt(s3, "TAK", LC + DW/2 + Inches(0.05), y + DH/2 - Inches(0.34),
    Inches(0.5), Inches(0.28), fs=Pt(9), bold=True, color=C_DECISION)
# NO branch ↓
arrow_d(s3, LC, y + DH, y + DH + Inches(0.3), C_SYSTEM)
txt(s3, "NIE", LC + Inches(0.08), y + DH + Inches(0.02),
    Inches(0.5), Inches(0.25), fs=Pt(9), bold=True, color=C_SYSTEM)
y += DH + Inches(0.3)

# Auth check
diamond(s3, LC - DW/2, y, DW, DH, RGBColor(0x2A, 0x10, 0x00),
        "from_id ==\nCHAT_ID?", fs=Pt(11), bold=True,
        fc=C_DECISION, line=C_DECISION)
# NO → skip
arrow_r(s3, LC + DW/2, LC + DW/2 + Inches(0.8), y + DH/2, C_DECISION)
rrect(s3, LC + DW/2 + Inches(0.8), y + DH/2 - Inches(0.22),
      Inches(1.35), Inches(0.44),
      RGBColor(0x28, 0x06, 0x06), "⛔  skip\n+update offset",
      fs=Pt(9), fc=C_DECISION, line=C_DECISION)
txt(s3, "NIE", LC + DW/2 + Inches(0.05), y + DH/2 - Inches(0.34),
    Inches(0.5), Inches(0.28), fs=Pt(9), bold=True, color=C_DECISION)
arrow_d(s3, LC, y + DH, y + DH + Inches(0.3), C_SYSTEM)
txt(s3, "TAK", LC + Inches(0.08), y + DH + Inches(0.02),
    Inches(0.5), Inches(0.25), fs=Pt(9), bold=True, color=C_SYSTEM)
y += DH + Inches(0.3)

# Router
rrect(s3, LC - NW/2, y, NW, NH, RGBColor(0x06, 0x24, 0x42),
      "🔀  COMMAND ROUTER\nparsuje tekst / callback",
      fs=Pt(10), bold=True, fc=C_TELEGRAM, line=C_TELEGRAM)
y_router_bot = y + NH

# Branch lines from router
branch_y = y + NH + Inches(0.25)
arrow_d(s3, LC, y + NH, branch_y, C_TELEGRAM)

# Horizontal branch bar
hline(s3, LC - Inches(1.5), LC + Inches(1.5), branch_y, C_TELEGRAM)

# Branch nodes below bar
branches = [
    (LC - Inches(1.5), "/status",         "systemctl\n+ gh pr list",   C_SYSTEM),
    (LC - Inches(0.5), "cancellation?",   "rm .pending*\nAnulowano",   C_DECISION),
    (LC + Inches(0.5), "/menu /start",    "inline\nkeyboard",          C_TELEGRAM),
    (LC + Inches(1.5), "/debate",         "→ debate.sh",               C_FILE),
]
for bx, label, sub, col in branches:
    vline(s3, bx, branch_y, branch_y + Inches(0.3), col)
    rrect(s3, bx - Inches(0.65), branch_y + Inches(0.3),
          Inches(1.3), Inches(0.62),
          CARD, f"{label}\n{sub}",
          fs=Pt(9), fc=col, line=col)

# Main path continues down center: confirmation → plan
arrow_d(s3, LC, branch_y + Inches(0.92), branch_y + Inches(1.2), C_TELEGRAM)
y = branch_y + Inches(1.2)

diamond(s3, LC - DW/2, y, DW, DH, RGBColor(0x20, 0x10, 0x00),
        ".pending_action\n+ confirmation?", fs=Pt(10), bold=True,
        fc=C_FILE, line=C_FILE)
arrow_d(s3, LC, y + DH, y + DH + Inches(0.3), C_SYSTEM)
txt(s3, "NIE → PLAN", LC + Inches(0.08), y + DH + Inches(0.01),
    Inches(1.1), Inches(0.28), fs=Pt(9), bold=True, color=C_SYSTEM)
# YES → EXEC
arrow_r(s3, LC + DW/2, LC + DW/2 + Inches(0.6), y + DH/2, C_EXEC)
rrect(s3, LC + DW/2 + Inches(0.6), y + DH/2 - Inches(0.22),
      Inches(1.4), Inches(0.44),
      RGBColor(0x28, 0x06, 0x06), "→ EXEC PHASE\n(prawa kolumna)",
      fs=Pt(9), fc=C_EXEC, line=C_EXEC)
txt(s3, "TAK", LC + DW/2 + Inches(0.05), y + DH/2 - Inches(0.34),
    Inches(0.5), Inches(0.28), fs=Pt(9), bold=True, color=C_EXEC)

# ── RIGHT COLUMN: Plan Phase & Exec Phase ───────────────────────
# PLAN PHASE
ry = Inches(1.5)
rrect(s3, RC - NW/2, ry, NW, Inches(0.82),
      RGBColor(0x10, 0x06, 0x2A),
      "📥  PLAN PHASE\nget_context():\nopen PRs · agents.md · git status",
      fs=Pt(10), fc=C_CLAUDE, line=C_CLAUDE)
ry += Inches(0.82)
arrow_d(s3, RC, ry, ry + Inches(0.28), C_CLAUDE)
ry += Inches(0.28)

rrect(s3, RC - NW/2, ry, NW, NH,
      RGBColor(0x18, 0x06, 0x38),
      "🤖  claude -p\n--model sonnet-4-6  timeout 120s",
      fs=Pt(10), fc=C_CLAUDE, line=C_CLAUDE)
ry += NH
arrow_d(s3, RC, ry, ry + Inches(0.28), C_CLAUDE)
ry += Inches(0.28)

diamond(s3, RC - DW/2, ry, DW, DH,
        RGBColor(0x1A, 0x08, 0x30),
        "response ends\nwith CONFIRM?", fs=Pt(11), bold=True,
        fc=C_CLAUDE, line=C_CLAUDE)
# YES branch →
arrow_r(s3, RC + DW/2, RC + DW/2 + Inches(0.55), ry + DH/2, C_FILE)
rrect(s3, RC + DW/2 + Inches(0.55), ry + DH/2 - Inches(0.33),
      Inches(1.55), Inches(0.66),
      RGBColor(0x2A, 0x1A, 0x00),
      "save\n.pending_action\nshow buttons",
      fs=Pt(9), fc=C_FILE, line=C_FILE)
txt(s3, "TAK", RC + DW/2 + Inches(0.05), ry + DH/2 - Inches(0.35),
    Inches(0.5), Inches(0.28), fs=Pt(9), bold=True, color=C_FILE)
# NO → send plain
arrow_d(s3, RC, ry + DH, ry + DH + Inches(0.25), C_CLAUDE)
txt(s3, "NIE", RC + Inches(0.08), ry + DH + Inches(0.01),
    Inches(0.5), Inches(0.25), fs=Pt(9), bold=True, color=C_CLAUDE)
ry += DH + Inches(0.25)

rrect(s3, RC - NW/2, ry, NW, NH,
      RGBColor(0x06, 0x28, 0x50),
      "📤  tg_send(response)\nprosta odpowiedź tekst",
      fs=Pt(10), fc=C_TELEGRAM, line=C_TELEGRAM)
ry += NH + Inches(0.5)

# EXEC PHASE
rect(s3, RC - NW/2 - Inches(0.1), ry - Inches(0.1),
     NW + Inches(0.2), Inches(2.25), RGBColor(0x1E, 0x05, 0x05),
     line=C_EXEC, lw=Pt(1.0))
txt(s3, "EXEC PHASE", RC - NW/2, ry - Inches(0.05), NW, Inches(0.28),
    fs=Pt(10), bold=True, color=C_EXEC, align=PP_ALIGN.CENTER)

rrect(s3, RC - NW/2, ry + Inches(0.3), NW, NH,
      RGBColor(0x28, 0x06, 0x06),
      "🤖  claude -p --allowedTools Bash(*)\ntimeout 300s  model: default",
      fs=Pt(10), fc=C_EXEC, line=C_EXEC)
arrow_d(s3, RC, ry + Inches(0.3) + NH, ry + Inches(0.3) + NH + Inches(0.25), C_EXEC)

rrect(s3, RC - NW/2, ry + Inches(0.3) + NH + Inches(0.25), NW, NH,
      RGBColor(0x10, 0x10, 0x10),
      "⚙️  Bash: git · gh · curl\nna serwerze rafserver",
      fs=Pt(10), fc=C_GITHUB, line=C_GITHUB)
arrow_d(s3, RC, ry + Inches(0.3) + NH * 2 + Inches(0.25),
        ry + Inches(0.3) + NH * 2 + Inches(0.5), C_TELEGRAM)

rrect(s3, RC - NW/2, ry + Inches(0.3) + NH * 2 + Inches(0.5), NW, NH,
      RGBColor(0x06, 0x28, 0x50),
      "📤  tg_send(wynik)\n+ update offset",
      fs=Pt(10), fc=C_TELEGRAM, line=C_TELEGRAM)

slide_num(s3, 3)


# ════════════════════════════════════════════════════════════════════
# SLIDE 4 — DEBATE PIPELINE
# ════════════════════════════════════════════════════════════════════

s4 = prs.slides.add_slide(BLANK)
bg(s4)
header(s4, "Debate Pipeline — /debate <temat>", "Claude ↔ Gemini · max 6 rund · consensus detection", C_FILE)

# Top: trigger
OW = Inches(2.5)
OH = Inches(0.65)
CX = Inches(6.67)

y = Inches(1.5)
rrect(s4, CX - OW/2, y, OW, OH,
      RGBColor(0x0A, 0x20, 0x3A),
      "📩  /debate <temat>\nlub przycisk w menu",
      fs=Pt(10), bold=True, fc=C_TELEGRAM, line=C_TELEGRAM)
y += OH
arrow_d(s4, CX, y, y + Inches(0.28), C_FILE)
y += Inches(0.28)

# Load context
rrect(s4, CX - OW/2, y, OW, OH,
      RGBColor(0x20, 0x16, 0x00),
      "📦  load context\ncode · PRs · TASKS.md · agents.md",
      fs=Pt(10), fc=C_FILE, line=C_FILE)
y += OH
arrow_d(s4, CX, y, y + Inches(0.25), C_FILE)
y += Inches(0.25)

# Loop box
loop_y = y
loop_h = Inches(2.55)
rect(s4, CX - Inches(5.2), y, Inches(10.4), loop_h,
     RGBColor(0x0C, 0x1C, 0x0C), line=C_SYSTEM, lw=Pt(1.2))
txt(s4, "🔁  RUNDA (max 6)", CX - Inches(5.2) + Inches(0.15), y + Inches(0.05),
    Inches(2.5), Inches(0.3), fs=Pt(9), bold=True, color=C_SYSTEM)

# Inside loop: Claude → Gemini → Haiku check
inner_y = y + Inches(0.4)
node_h = Inches(0.65)
node_w = Inches(2.8)

# Claude node
claude_cx = CX - Inches(3.0)
rrect(s4, claude_cx - node_w/2, inner_y, node_w, node_h,
      RGBColor(0x22, 0x08, 0x3A),
      "🤖  Claude\nsonnet-4-6\nrola: security dev",
      fs=Pt(10), fc=C_CLAUDE, line=C_CLAUDE)

# Gemini node
gemini_cx = CX + Inches(3.0)
rrect(s4, gemini_cx - node_w/2, inner_y, node_w, node_h,
      RGBColor(0x06, 0x22, 0x12),
      "✨  Gemini\ngemini -p\nrola: code reviewer",
      fs=Pt(10), fc=C_SYSTEM, line=C_SYSTEM)

# Arrow Claude → Gemini
arrow_r(s4, claude_cx + node_w/2, gemini_cx - node_w/2, inner_y + node_h/2, C_CLAUDE)
txt(s4, "powiedział Claude:", CX - Inches(1.0), inner_y + node_h/2 - Inches(0.30),
    Inches(2.0), Inches(0.28), fs=Pt(8), color=C_CLAUDE, align=PP_ALIGN.CENTER)

# Arrow Gemini response back to loop / consensus
arrow_d(s4, gemini_cx, inner_y + node_h, inner_y + node_h + Inches(0.3), C_SYSTEM)

# Consensus check
cons_y = inner_y + node_h + Inches(0.3)
cons_cx = CX + Inches(3.0)
DW2 = Inches(2.8)
DH2 = Inches(0.75)
diamond(s4, cons_cx - DW2/2, cons_y, DW2, DH2,
        RGBColor(0x12, 0x1A, 0x06),
        "Haiku:\nCONSENSUS?", fs=Pt(10), bold=True,
        fc=C_SYSTEM, line=C_SYSTEM)

# CONSENSUS → break (right)
arrow_r(s4, cons_cx + DW2/2, cons_cx + DW2/2 + Inches(1.0),
        cons_y + DH2/2, C_SYSTEM)
rrect(s4, cons_cx + DW2/2 + Inches(1.0), cons_y + DH2/2 - Inches(0.22),
      Inches(1.3), Inches(0.44),
      RGBColor(0x06, 0x22, 0x06), "✅  BREAK\ndone", fs=Pt(9), fc=C_SYSTEM, line=C_SYSTEM)
txt(s4, "CONSENSUS", cons_cx + DW2/2 + Inches(0.03), cons_y + DH2/2 - Inches(0.35),
    Inches(0.95), Inches(0.28), fs=Pt(8), bold=True, color=C_SYSTEM)

# CONTINUE → loop back up to Claude
arrow_l(s4, cons_cx - DW2/2, CX - Inches(5.0), cons_y + DH2/2, C_FILE)
vline(s4, CX - Inches(5.0), inner_y + node_h/2, cons_y + DH2/2, C_FILE)
arrow_r(s4, CX - Inches(5.0), claude_cx - node_w/2, inner_y + node_h/2, C_FILE)
txt(s4, "CONTINUE", CX - Inches(4.8), cons_y + DH2/2 + Inches(0.04),
    Inches(0.9), Inches(0.25), fs=Pt(8), bold=True, color=C_FILE)

y = loop_y + loop_h
arrow_d(s4, CX, y, y + Inches(0.25), C_CLAUDE)
y += Inches(0.25)

# Summary
rrect(s4, CX - OW/2, y, OW, OH,
      RGBColor(0x18, 0x08, 0x30),
      "📝  claude summary\nkluczowe ustalenia + TASK: ekstrakcja",
      fs=Pt(10), fc=C_CLAUDE, line=C_CLAUDE)
y += OH
arrow_d(s4, CX, y, y + Inches(0.25), C_CLAUDE)
y += Inches(0.25)

diamond(s4, CX - Inches(1.5), y, Inches(3.0), DH2,
        RGBColor(0x20, 0x12, 0x00),
        "TASK: lines\nnaleziono?", fs=Pt(11), bold=True,
        fc=C_FILE, line=C_FILE)
# TAK →
arrow_r(s4, CX + Inches(1.5), CX + Inches(2.0), y + DH2/2, C_FILE)
rrect(s4, CX + Inches(2.0), y + DH2/2 - Inches(0.45),
      Inches(2.9), Inches(0.9),
      RGBColor(0x10, 0x22, 0x10),
      "🤖  buttons:\n\"Deleguj do Jules\"\njules new --repo",
      fs=Pt(9), fc=C_SYSTEM, line=C_SYSTEM)
txt(s4, "TAK", CX + Inches(1.55), y + DH2/2 - Inches(0.36),
    Inches(0.5), Inches(0.25), fs=Pt(9), bold=True, color=C_FILE)
# NIE ↓
arrow_d(s4, CX, y + DH2, y + DH2 + Inches(0.25), C_TELEGRAM)
txt(s4, "NIE", CX + Inches(0.08), y + DH2 + Inches(0.01),
    Inches(0.5), Inches(0.25), fs=Pt(9), bold=True, color=C_TELEGRAM)
y += DH2 + Inches(0.25)

rrect(s4, CX - OW/2, y, OW, OH,
      RGBColor(0x06, 0x22, 0x40),
      "📤  tg_send(summary)\ndo Telegrama",
      fs=Pt(10), fc=C_TELEGRAM, line=C_TELEGRAM)

slide_num(s4, 4)


# ════════════════════════════════════════════════════════════════════
# SLIDE 5 — PR REVIEW PIPELINE
# ════════════════════════════════════════════════════════════════════

s5 = prs.slides.add_slide(BLANK)
bg(s5)
header(s5, "PR Review Pipeline — Auto", "cron co 5 minut · 1 PR per run · .reviewed_prs state", C_CLAUDE)

# Vertical centered flow
FC = Inches(3.5)   # left fork center
RC5 = Inches(9.5)  # right side for detail
NW5 = Inches(3.4)
NH5 = Inches(0.68)

y = Inches(1.5)

# Cron trigger
oval(s5, FC - Inches(1.5), y, Inches(3.0), Inches(0.52),
     C_SYSTEM, "⏰  cron  — każde 5 min", fs=Pt(12), bold=True, fc=TEXT_DARK)
y += Inches(0.52)
arrow_d(s5, FC, y, y + Inches(0.28), C_SYSTEM)
y += Inches(0.28)

# gh pr list
rrect(s5, FC - NW5/2, y, NW5, NH5,
      RGBColor(0x10, 0x10, 0x10),
      "🐙  gh pr list --repo mrBTL/VulnApp-Light\n--state open  →  JSON: number · title · author · date",
      fs=Pt(10), fc=C_GITHUB, line=C_GITHUB)
y += NH5
arrow_d(s5, FC, y, y + Inches(0.25), C_GITHUB)
y += Inches(0.25)

# Any PRs?
diamond(s5, FC - Inches(1.6), y, Inches(3.2), Inches(0.78),
        RGBColor(0x28, 0x06, 0x06),
        "open PRs\nexistują?", fs=Pt(11), bold=True,
        fc=C_DECISION, line=C_DECISION)
# NO → exit
arrow_r(s5, FC + Inches(1.6), FC + Inches(2.2), y + Inches(0.39), C_DECISION)
oval(s5, FC + Inches(2.2), y + Inches(0.15), Inches(1.5), Inches(0.48),
     RGBColor(0x20, 0x06, 0x06), "exit(0)", fs=Pt(11), bold=True, fc=C_DECISION)
txt(s5, "NIE", FC + Inches(1.65), y + Inches(0.05),
    Inches(0.5), Inches(0.25), fs=Pt(9), bold=True, color=C_DECISION)
# YES ↓
arrow_d(s5, FC, y + Inches(0.78), y + Inches(0.78) + Inches(0.25), C_SYSTEM)
txt(s5, "TAK", FC + Inches(0.08), y + Inches(0.78) + Inches(0.01),
    Inches(0.5), Inches(0.25), fs=Pt(9), bold=True, color=C_SYSTEM)
y += Inches(0.78) + Inches(0.25)

# For each PR
rrect(s5, FC - NW5/2, y, NW5, NH5,
      RGBColor(0x06, 0x20, 0x08),
      "🔁  foreach PR",
      fs=Pt(12), bold=True, fc=C_SYSTEM, line=C_SYSTEM)
y += NH5
arrow_d(s5, FC, y, y + Inches(0.25), C_SYSTEM)
y += Inches(0.25)

# Already reviewed?
diamond(s5, FC - Inches(1.6), y, Inches(3.2), Inches(0.78),
        RGBColor(0x1E, 0x14, 0x00),
        "PR# w\n.reviewed_prs?", fs=Pt(11), bold=True,
        fc=C_FILE, line=C_FILE)
# YES → skip
arrow_r(s5, FC + Inches(1.6), FC + Inches(2.2), y + Inches(0.39), C_FILE)
rrect(s5, FC + Inches(2.2), y + Inches(0.15), Inches(1.5), Inches(0.48),
      RGBColor(0x20, 0x14, 0x00), "skip\n(już zrecenzowany)",
      fs=Pt(9), fc=C_FILE, line=C_FILE)
txt(s5, "TAK", FC + Inches(1.65), y + Inches(0.05),
    Inches(0.5), Inches(0.25), fs=Pt(9), bold=True, color=C_FILE)
# NO ↓
arrow_d(s5, FC, y + Inches(0.78), y + Inches(0.78) + Inches(0.25), C_SYSTEM)
txt(s5, "NIE", FC + Inches(0.08), y + Inches(0.78) + Inches(0.01),
    Inches(0.5), Inches(0.25), fs=Pt(9), bold=True, color=C_SYSTEM)
y += Inches(0.78) + Inches(0.25)

# Get diff
rrect(s5, FC - NW5/2, y, NW5, NH5,
      RGBColor(0x10, 0x10, 0x10),
      "📄  gh pr diff <number>\nhead -150  (truncated diff)",
      fs=Pt(10), fc=C_GITHUB, line=C_GITHUB)
y += NH5
arrow_d(s5, FC, y, y + Inches(0.25), C_CLAUDE)
y += Inches(0.25)

# Claude review
rrect(s5, FC - NW5/2, y, NW5, NH5,
      RGBColor(0x16, 0x06, 0x2C),
      "🤖  claude -p (no tools)\nreview: APPROVE / CHANGES NEEDED / CRITICAL\nmax 300 chars · agents.md context",
      fs=Pt(10), fc=C_CLAUDE, line=C_CLAUDE)
y += NH5
arrow_d(s5, FC, y, y + Inches(0.25), C_TELEGRAM)
y += Inches(0.25)

# Send to Telegram
rrect(s5, FC - NW5/2, y, NW5, NH5,
      RGBColor(0x06, 0x24, 0x44),
      "📤  tg_send: PR# · title · review\n/merge N  /close N  /comment N <tekst>",
      fs=Pt(10), fc=C_TELEGRAM, line=C_TELEGRAM)
y += NH5
arrow_d(s5, FC, y, y + Inches(0.25), C_FILE)
y += Inches(0.25)

# Save state
rrect(s5, FC - NW5/2, y, NW5, NH5 * 0.8,
      RGBColor(0x22, 0x14, 0x00),
      "💾  echo PR# >> .reviewed_prs  →  exit(0)",
      fs=Pt(10), fc=C_FILE, line=C_FILE)

# ── RIGHT SIDE: Telegram message anatomy ────────────────────────
txt(s5, "Wiadomość w Telegram:", RC5, Inches(1.5), Inches(3.5), Inches(0.35),
    fs=Pt(11), bold=True, color=C_TELEGRAM)

msg_lines = [
    ("🤖 Jules PR #42",                     C_TELEGRAM),
    ("📝  Fix: SQL injection in search",    TEXT_W),
    ("",                                     TEXT_W),
    ("🔍 Claude review:",                   C_CLAUDE),
    ("CHANGES NEEDED — query not",          TEXT_W),
    ("parameterized, line 87",              TEXT_W),
    ("",                                     TEXT_W),
    ("✅  /merge 42",                        C_SYSTEM),
    ("❌  /close 42",                       C_DECISION),
    ("✏️  /comment 42 <tekst>",            C_FILE),
]
msg_y = Inches(1.9)
msg_bg = rect(s5, RC5 - Inches(0.1), msg_y - Inches(0.1),
              Inches(3.6), Inches(3.2), CARD, line=C_TELEGRAM)
for line, col in msg_lines:
    txt(s5, line, RC5 + Inches(0.05), msg_y, Inches(3.4), Inches(0.3),
        fs=Pt(10), color=col)
    msg_y += Inches(0.28)

# State file diagram
txt(s5, ".reviewed_prs (state):", RC5, Inches(5.3), Inches(3.5), Inches(0.35),
    fs=Pt(11), bold=True, color=C_FILE)
state_bg = rect(s5, RC5 - Inches(0.1), Inches(5.65), Inches(3.6), Inches(1.2),
               CARD, line=C_FILE)
for i, line in enumerate(["37", "38", "39", "40", "41", "42  ← nowy"]):
    col = C_FILE if i == 5 else TEXT_DIM
    bold_ = i == 5
    txt(s5, line, RC5 + Inches(0.1), Inches(5.75) + i * Inches(0.17),
        Inches(2.0), Inches(0.2), fs=Pt(10), color=col, bold=bold_)

slide_num(s5, 5)


# ════════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════════
out = "Jules_Bot_Architecture.pptx"
prs.save(out)
print(f"Saved: {out}")
